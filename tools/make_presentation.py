#!/usr/bin/env python3
"""Build a PowerPoint presentation summarising today's block-offset research."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RESEARCH = Path("/scratch/code/debris/research")
OUT      = RESEARCH / "block_offset_exploration.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x12, 0x12, 0x18)
HEADING = RGBColor(0xFF, 0xB3, 0x00)   # amber
BODY    = RGBColor(0xDD, 0xDD, 0xDD)
DIM     = RGBColor(0x88, 0x88, 0x99)
ACCENT  = RGBColor(0x00, 0x99, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def _bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def _tb(slide, text, l, t, w, h, size=18, bold=False,
        color=BODY, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.color.rgb = color
    return txb

def _img(slide, path, l, t, w, h=None):
    if h is None:
        slide.shapes.add_picture(str(path), l, t, width=w)
    else:
        slide.shapes.add_picture(str(path), l, t, width=w, height=h)

def _rule(slide, t, color=HEADING, thickness=Pt(1.5)):
    from pptx.util import Emu
    line = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.LINE → use freeform connector instead
        Inches(0.3), t, Inches(12.73), Emu(0),
    )
    line.line.color.rgb = color
    line.line.width     = thickness

def _title_slide(title, subtitle=None):
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _tb(slide, title,
        Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.8),
        size=40, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
    if subtitle:
        _tb(slide, subtitle,
            Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.2),
            size=22, color=BODY, align=PP_ALIGN.CENTER)
    return slide

def _section_slide(title, body_lines):
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _tb(slide, title,
        Inches(0.4), Inches(0.25), Inches(12.53), Inches(0.65),
        size=26, bold=True, color=HEADING)
    _rule(slide, Inches(1.0))
    y = Inches(1.15)
    for line in body_lines:
        size = 15 if line.startswith("  ") else 17
        col  = DIM  if line.startswith("  ") else BODY
        _tb(slide, line.strip(),
            Inches(0.5), y, Inches(12.3), Inches(0.55),
            size=size, color=col)
        y += Inches(0.52)
    return slide

def _img_slide(title, img_path, caption=None, top_frac=0.88):
    """Full-width image slide."""
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _tb(slide, title,
        Inches(0.4), Inches(0.15), Inches(12.53), Inches(0.55),
        size=22, bold=True, color=HEADING)
    img_h = H * top_frac - Inches(0.75)
    _img(slide, img_path, Inches(0.2), Inches(0.75), Inches(12.93), img_h)
    if caption:
        _tb(slide, caption,
            Inches(0.4), H - Inches(0.55), Inches(12.53), Inches(0.45),
            size=12, color=DIM, align=PP_ALIGN.CENTER)
    return slide

def _two_col(title, left_lines, img_path, img_caption=None):
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _tb(slide, title,
        Inches(0.4), Inches(0.15), Inches(12.53), Inches(0.55),
        size=22, bold=True, color=HEADING)
    _rule(slide, Inches(0.8))
    y = Inches(0.95)
    for line in left_lines:
        size = 14 if line.startswith("  ") else 16
        col  = DIM  if line.startswith("  ") else BODY
        _tb(slide, line.strip(),
            Inches(0.4), y, Inches(5.0), Inches(0.52),
            size=size, color=col, wrap=True)
        y += Inches(0.50)
    _img(slide, img_path, Inches(5.6), Inches(0.75), Inches(7.5))
    if img_caption:
        _tb(slide, img_caption,
            Inches(5.6), H - Inches(0.5), Inches(7.5), Inches(0.4),
            size=11, color=DIM, align=PP_ALIGN.CENTER)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════════════

# 1. Title
_title_slide(
    "Bitcoin Power Law Floor Model\nBlock-Space Exploration",
    "Research session  ·  April 2026",
)

# 2. Agenda
_section_slide("Today's agenda", [
    "1.  Background: why model price in block-space?",
    "2.  BM floor model sweep — bottom fraction × fit quantile × weight mode",
    "3.  Band model — fit QN% to ±5-percentile residual band",
    "4.  Temporal approaches to finding the optimal block origin",
    "    A2: temporal quantile consistency (bin percentiles)",
    "    A1: time-stratified floor selection",
    "5.  q*(offset) — the convergence-quantile method",
    "6.  Key findings & open questions",
])

# 3. Background
_section_slide("Background: block-space power law", [
    "Model:  log₁₀(price) = a + b · log₁₀(block − offset)",
    "",
    "• 'block' = blockheight at midnight UTC for each daily price",
    "• 'offset' = the true time origin  (genesis block is 0, but optimal may differ)",
    "• In log-log space the model is linear — a straight line",
    "• The exponent b ≈ 5–6 for most fits to date",
    "",
    "Key question explored today:",
    "  What offset produces the most stable, temporally-uniform floor fit?",
])

# 4. BM floor model recap
_two_col(
    "BM floor model: parameter sweep",
    [
        "Two tuning knobs:",
        "  • Bottom fraction: keep lowest N% of data",
        "    (selected by QR residuals)",
        "  • Fit quantile: fit QM% to that subset",
        "",
        "Swept across:",
        "  • Fractions: 5%, 20%",
        "  • Quantiles: 5%, 20%",
        "  • 4 weight modes",
        "  • Offsets 0 → 37,500 blocks (step 150)",
        "",
        "Problem identified:",
        "  Global percentile selection is biased toward",
        "  early data — Q5/10 miss all post-500k blocks.",
    ],
    RESEARCH / "blocksweep.jpg",
    "R² and exponent vs block offset  |  4 weight modes  |  BM floor + band models",
)

# 5. Weight modes explained
_section_slide("Weight modes", [
    "Unweighted       — all data points equal",
    "",
    "1/t weighted     — emphasise early data",
    "  Early price history carries more 'independent' information per block",
    "  (fewer blocks per halving cycle)",
    "",
    "1/√t weighted    — softer early emphasis, compromise between the two",
    "",
    "Log-density weighted  — 1/KDE(log t)",
    "  Corrects for uneven sampling density in log-time.",
    "  Gives each log-time era equal effective weight regardless of",
    "  how many calendar days fall in that bin.",
    "",
    "Key finding: at offset ≈ 23k–26k blocks, all four weight modes",
    "  produce near-identical slopes (~5.05–5.13).",
])

# 6. Band model
_section_slide("Band model — new approach", [
    "Motivation: BM floor always takes the bottom tail → ignores recent data.",
    "",
    "Band model for target percentile N:",
    "  1.  Fit reference QR at Q N% to full data → residuals",
    "  2.  Keep data in the [(N−5)th, (N+5)th] percentile band of residuals",
    "      → a symmetric 10%-wide slice centred on the Nth percentile",
    "  3.  Fit QR at Q N% to that slice",
    "",
    "Tested for N = 5%, 10%, 25%",
    "",
    "Advantage: selects data near the line of interest from every era,",
    "  not just the earliest bear-market lows.",
    "",
    "Finding: Q25% band gives the most stable slope across weight modes.",
])

# 7. Band model fits offset=0
_img_slide(
    "Band model fits — offset = 0",
    RESEARCH / "bandmodel_fits.jpg",
    "Colored scatter = data in each band  |  Lines = fitted QR  |  Left: log-log  |  Right: log-linear",
)

# 8. Band model fits offset=18750
_img_slide(
    "Band model fits — offset = 18,750 blocks",
    RESEARCH / "bandmodel_fits_offset18750.jpg",
    "Shifting origin compresses early log-time, making recent data more prominent in the fit",
)

# 9. Band model coefficients
_section_slide("Band model Q25%: slope & intercept comparison", [
    "log₁₀(price) = a + b · log₁₀(block − offset)",
    "",
    "Offset = 0:",
    "  Unweighted      a = −27.018   b = 5.320",
    "  1/t             a = −28.697   b = 5.616",
    "  1/√t            a = −28.042   b = 5.500",
    "  Log-density     a = −28.787   b = 5.632",
    "",
    "Offset = 18,750:",
    "  Unweighted      a = −25.372   b = 5.053",
    "  1/t             a = −25.827   b = 5.133",
    "  1/√t            a = −25.606   b = 5.094",
    "  Log-density     a = −25.714   b = 5.113",
    "",
    "At offset 18,750 the four weight modes converge tightly (b = 5.05–5.13).",
])

# 10. Temporal approaches — motivation
_section_slide("Problem: global selection is temporally biased", [
    "Observation: recent Bitcoin cycles do not reach as low (on a percentile basis)",
    "as early cycles. Each cycle's floor is higher relative to trend.",
    "",
    "Consequence:",
    "  • Global Q5% or Q10% QR selects almost exclusively early data",
    "  • Fitted line has a steep exponent driven by 2010–2013 lows",
    "  • No data after block ~500,000 falls below the Q5% line",
    "",
    "Two temporal approaches tested:",
    "",
    "  A2 — Temporal quantile consistency",
    "    Bin log-time into 10 equal bins → percentile per bin → WLS fit",
    "    Score: variance of per-bin fraction-below (lower = more uniform)",
    "",
    "  A1 — Time-stratified floor selection",
    "    Bin log-time → take bottom P% per bin → pool → global QR",
    "    Score: R² on pooled subset",
])

# 11. Temporal sweep
_img_slide(
    "Temporal sweep: A2 and A1 scores vs block offset",
    RESEARCH / "temporal_sweep.jpg",
    "Cols: A2 R², A2 consistency variance (↓ better), A1 R², A1 exponent  |  Rows: 4 weight modes",
)

# 12. Best temporal fits
_img_slide(
    "Best fits at optimal (Q-level, offset) per weight mode",
    RESEARCH / "temporal_fits.jpg",
    "Blue solid = A2 best (Q-level, offset) by min cvar  |  Red dashed = A1 best combo by max R²",
)

# 13. Convergence-quantile insight
_section_slide("The convergence-quantile insight", [
    "Observation: lower quantiles (Q5%, Q10%) → steeper slope (higher b)",
    "             higher quantiles (Q50%, Q75%) → flatter slope (lower b)",
    "",
    "These quantile lines fan out from an origin — they converge somewhere.",
    "",
    "Key insight:",
    "  For each offset, q*(offset) = the quantile whose A2 cvar is minimised",
    "  = the quantile most uniformly distributed across all time eras",
    "",
    "  If the block origin is correctly specified, the model should be",
    "  most consistent at Q50% (symmetric residuals across eras).",
    "",
    "Two targets:",
    "  q*(offset) = 50%  →  optimal offset, full-distribution model",
    "  q*(offset) = 25%  →  optimal offset, bear-market-only model",
    "",
    "Finding: the Q25% and Q50% crossings are nearly identical (~40–54 blocks apart).",
])

# 14. q* sweep
_img_slide(
    "q*(offset) — convergence-quantile sweep",
    RESEARCH / "qstar_sweep.jpg",
    "Top row: cvar heatmap (darker = more consistent) with q* overlaid  |  "
    "Bottom row: q*(offset) curves with Q25% and Q50% crossings annotated",
)

# 15. q* optimal fits
_img_slide(
    "q*-crossing fits: optimal (Q-level, offset) per weight mode",
    RESEARCH / "qstar_fits.jpg",
    "Lines fitted at the exact interpolated crossing offset  |  "
    "Red dashed = q*=25% target  |  Blue solid = q*=50% target  |  "
    "Legend shows offset, intercept (a), slope (b)",
)

# 16. Secondary feature observation
_section_slide("Observation: secondary feature at high quantiles / low offset", [
    "Inspecting the 3D cvar surface revealed a second region of low cvar:",
    "  • Offset ≈ 0–5,000 blocks",
    "  • Quantile ≈ Q75–85%",
    "",
    "This may represent a 'bull market line' — a power law that consistently",
    "tracks the upper portion of price distribution across all eras.",
    "",
    "Key question: is this feature stable (wide, shallow valley) or fragile",
    "(narrow spike sensitive to offset choice)?",
    "",
    "Investigation:",
    "  • Fine-resolution sweep: step 25 blocks (0–8k), step 50 (8k–37.5k)",
    "  • Q steps of 2.5% (vs 5% before)",
    "  • Normalised cvar slices to compare sharpness across Q levels",
    "  • Valley width = offset range where cvar ≤ 1.10 × cvar_min",
])

# 17. Fine 3D surface
_img_slide(
    "Fine-resolution cvar surface — log-density weighted",
    RESEARCH / "qstar_fine_3d.jpg",
    "Step 25 blocks (0–8k) and 50 blocks (8k–37.5k)  |  Q steps 2.5%  |  "
    "Yellow slice = Q80%  |  Two distinct low-cvar features visible",
)

# 18. cvar slices
_img_slide(
    "cvar vs offset at fixed quantile levels (normalised)",
    RESEARCH / "qstar_slices.jpg",
    "Left: low quantiles Q10–30% (main valley)  |  Right: high quantiles Q65–90% (secondary feature)\n"
    "Curves normalised to cvar_min = 1  |  Wider = more stable model  |  "
    "Dashed line = 10% above minimum",
)

# 19. Valley width analysis
_img_slide(
    "Valley width analysis across all quantile levels",
    RESEARCH / "qstar_valley_width.jpg",
    "Top: valley width (blocks where cvar ≤ 1.10 × min)  |  "
    "Middle: optimal offset per Q level  |  Bottom: best achievable cvar\n"
    "Blue shading = main valley region (Q5–35%)  |  Orange shading = secondary feature (Q65–90%)",
)

# 20. Key findings
_section_slide("Key findings", [
    "1.  Optimal block offset  ≈  23,000–26,000 blocks",
    "    (all weight modes and both target quantiles agree within ~3,000 blocks)",
    "",
    "2.  The Q25% and Q50% crossings are nearly identical (~40–54 blocks apart)",
    "    → robust to assumption about what portion of the distribution is predictable",
    "",
    "3.  At optimal offset, weight modes converge tightly (b ≈ 4.76–5.09)",
    "    → flatter than uncorrected fits at offset=0 (b ≈ 5.3–5.6)",
    "",
    "4.  A secondary low-cvar feature exists near offset ≈ 0–5k / Q75–85%",
    "    → candidate 'bull market line'; stability vs main valley TBD from valley-width plot",
    "",
    "5.  Band model (±5 pct band around QN%) gives stable, era-spanning fits",
    "    → better behaved than global BM floor for recent data",
])

# 21. Open questions
_section_slide("Open questions", [
    "• Is the optimal offset of ~25,000 blocks ≈ 172 days economically meaningful?",
    "  (Genesis block: 2009-01-03  →  offset of 25k ≈ mid-2009, before first dollar trades)",
    "",
    "• Should we refit the main site model (currently t₀ = 2009-07-25 / ~110k blocks)?",
    "  Calendar-space optimal date and block-space optimal offset are not the same thing.",
    "",
    "• Is the secondary high-Q / low-offset feature a genuine 'bull market line'?",
    "  Valley-width analysis will show if it is stable enough to be a useful model.",
    "",
    "• Can the band model be incorporated as a registered price model on the site?",
    "",
    "• Does the optimal offset shift as new price history accumulates?",
    "  (Run the sweep again after the next cycle bottom to test stability.)",
])

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
